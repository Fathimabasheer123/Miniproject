from flask import Blueprint, render_template, request, redirect, url_for, session, jsonify, send_file, flash
from werkzeug.utils import secure_filename
import PyPDF2
from pptx import Presentation
import json
import os
import re
from transformers import pipeline, AutoTokenizer
from sentence_transformers import SentenceTransformer, util
from nltk.corpus import wordnet as wn
import random
from collections import Counter
import io
import ssl
from datetime import datetime
from werkzeug.security import generate_password_hash, check_password_hash
import sqlite3
from functools import wraps

# NLTK setup with enhanced error handling
import nltk

try:
    _create_unverified_https_context = ssl._create_unverified_context
except AttributeError:
    pass
else:
    ssl._create_default_https_context = _create_unverified_https_context

# Try to download required NLTK data with fallback
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    try:
        nltk.download('punkt', quiet=True)
    except Exception:
        print("Warning: Could not download punkt tokenizer")

try:
    nltk.data.find('taggers/averaged_perceptron_tagger')
except LookupError:
    try:
        nltk.download('averaged_perceptron_tagger', quiet=True)
    except Exception:
        print("Warning: Could not download POS tagger")

try:
    nltk.data.find('corpora/wordnet')
except LookupError:
    try:
        nltk.download('wordnet', quiet=True)
    except Exception:
        print("Warning: Could not download WordNet. Distractors may be limited.")

# Create a Blueprint
routes = Blueprint("routes", __name__)

# Database setup
def get_db_connection():
    conn = sqlite3.connect('quizgen.db')
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    conn = get_db_connection()
    conn.execute('''
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE NOT NULL,
            email TEXT UNIQUE NOT NULL,
            password_hash TEXT NOT NULL,
            preferences TEXT DEFAULT '{"theme": "light"}',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    conn.execute('''
        CREATE TABLE IF NOT EXISTS quizzes (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            title TEXT NOT NULL,
            context TEXT,
            generated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (user_id) REFERENCES users (id)
        )
    ''')
    
    conn.execute('''
        CREATE TABLE IF NOT EXISTS questions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            quiz_id INTEGER NOT NULL,
            question_text TEXT NOT NULL,
            options TEXT NOT NULL,
            answer TEXT NOT NULL,
            FOREIGN KEY (quiz_id) REFERENCES quizzes (id)
        )
    ''')
    
    conn.execute('''
        CREATE TABLE IF NOT EXISTS quiz_attempts (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            quiz_id INTEGER NOT NULL,
            score INTEGER NOT NULL,
            total_questions INTEGER NOT NULL,
            answers TEXT NOT NULL,
            completed_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (user_id) REFERENCES users (id),
            FOREIGN KEY (quiz_id) REFERENCES quizzes (id)
        )
    ''')
    
    conn.commit()
    conn.close()

# Initialize database
init_db()

# Login required decorator
def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'user_id' not in session:
            return redirect(url_for('routes.login'))
        return f(*args, **kwargs)
    return decorated_function

# ---------------- Validation Helpers ----------------
def is_valid_username(username):
    return bool(re.match(r"^[A-Za-z0-9_]{3,}$", username))

def is_valid_email(email):
    return bool(re.match(r"^[\w\.-]+@[\w\.-]+\.\w+$", email))

def is_valid_password(password):
    if len(password) < 8:
        return False
    if not re.search(r"[A-Z]", password):
        return False
    if not re.search(r"[a-z]", password):
        return False
    if not re.search(r"[0-9]", password):
        return False
    if not re.search(r"[!@#$%^&*(),.?\":{}|<>]", password):
        return False
    return True

# ---------------- INDEX PAGE ----------------
@routes.route("/")
def index():
    return render_template("index.html")

# ---------------- LOGIN PAGE ----------------
@routes.route("/login", methods=["GET", "POST"])
def login():
    error = None
    if request.method == "POST":
        login_id = request.form.get("email")
        password = request.form.get("password")

        if not login_id or not password:
            error = "Both fields are required"
        else:
            conn = get_db_connection()
            user = conn.execute(
                'SELECT * FROM users WHERE username = ? OR email = ?', 
                (login_id, login_id)
            ).fetchone()
            conn.close()

            if user and check_password_hash(user['password_hash'], password):
                session['user_id'] = user['id']
                session['user'] = user['username']
                session['user_email'] = user['email']
                return redirect(url_for("routes.dashboard"))
            else:
                error = "Invalid login credentials"

    return render_template("login.html", error=error)

# ---------------- REGISTER PAGE ----------------
@routes.route("/register", methods=["GET", "POST"])
def register():
    error = None
    if request.method == "POST":
        username = request.form.get("username")
        email = request.form.get("email")
        password = request.form.get("password")

        if not username or not email or not password:
            error = "All fields are required"
        elif not is_valid_username(username):
            error = "Username must be at least 3 characters"
        elif not is_valid_email(email):
            error = "Invalid email format"
        elif not is_valid_password(password):
            error = "Password must be 8+ chars with uppercase, lowercase, digit, special char"
        else:
            conn = get_db_connection()
            
            # Check if email or username already exists
            existing_user = conn.execute(
                'SELECT * FROM users WHERE username = ? OR email = ?', 
                (username, email)
            ).fetchone()
            
            if existing_user:
                error = "Username or email already registered"
            else:
                # Create new user
                password_hash = generate_password_hash(password)
                cursor = conn.cursor()
                cursor.execute(
                    'INSERT INTO users (username, email, password_hash) VALUES (?, ?, ?)',
                    (username, email, password_hash)
                )
                conn.commit()
                
                # Get the new user
                user = conn.execute(
                    'SELECT * FROM users WHERE username = ?', (username,)
                ).fetchone()
                conn.close()
                
                # Log the user in
                session['user_id'] = user['id']
                session['user'] = user['username']
                session['user_email'] = user['email']
                return redirect(url_for("routes.dashboard"))

    return render_template("register.html", error=error)

# ---------------- DASHBOARD ----------------
@routes.route("/dashboard")
@login_required
def dashboard():
    conn = get_db_connection()
    
    # Get user stats
    user_id = session['user_id']
    
    # Quizzes completed
    quizzes_completed = conn.execute(
        'SELECT COUNT(*) FROM quiz_attempts WHERE user_id = ?', (user_id,)
    ).fetchone()[0]
    
    # Average score
    avg_score_result = conn.execute(
        'SELECT AVG(score * 100.0 / total_questions) FROM quiz_attempts WHERE user_id = ?', (user_id,)
    ).fetchone()[0]
    average_score = round(avg_score_result, 1) if avg_score_result else 0
    
    # Quizzes created
    quizzes_created = conn.execute(
        'SELECT COUNT(*) FROM quizzes WHERE user_id = ?', (user_id,)
    ).fetchone()[0]
    
    # Recent quizzes (last 5 attempts)
    recent_attempts = conn.execute('''
        SELECT q.title, qa.score, qa.total_questions, qa.completed_at 
        FROM quiz_attempts qa 
        JOIN quizzes q ON qa.quiz_id = q.id 
        WHERE qa.user_id = ? 
        ORDER BY qa.completed_at DESC 
        LIMIT 5
    ''', (user_id,)).fetchall()
    
    # Format recent quizzes for display
    recent_quizzes = []
    for attempt in recent_attempts:
        recent_quizzes.append({
            'title': attempt['title'],
            'score': round((attempt['score'] / attempt['total_questions']) * 100),
            'date': attempt['completed_at'][:10]  # Just the date part
        })
    
    # Calculate progress percentage (quizzes completed vs goal of 20 quizzes)
    quizzes_goal = 20
    progress_percentage = min(100, round((quizzes_completed / quizzes_goal) * 100))
    
    # Calculate days streak (placeholder implementation)
    # In a real app, you would calculate this based on consecutive days with quiz activity
    days_streak = random.randint(1, 10)  # Placeholder
    
    # Check achievements
    achievements = {
        'first_quiz': quizzes_completed >= 1,
        'quiz_master': quizzes_completed >= 10,
        'perfect_score': any(attempt['score'] == attempt['total_questions'] for attempt in recent_attempts),
        'streak_7': days_streak >= 7,
        'earned': 0,
        'total': 4
    }
    
    # Count earned achievements
    achievements['earned'] = sum(1 for k, v in achievements.items() if v and k not in ['earned', 'total'])
    
    conn.close()
    
    # Prepare stats for template
    stats = {
        'quizzes_completed': quizzes_completed,
        'average_score': average_score,
        'quizzes_created': quizzes_created,
        'progress_percentage': progress_percentage,
        'quizzes_goal': quizzes_goal,
        'days_streak': days_streak,
        'achievements': achievements
    }
    
    return render_template("dashboard.html", user=session['user'], stats=stats, recent_quizzes=recent_quizzes)

# ---------------- HISTORY PAGE ----------------
@routes.route("/history")
@login_required
def history():
    conn = get_db_connection()
    
    # Get quiz attempts for the user
    attempts = conn.execute('''
        SELECT qa.*, q.title, q.generated_at 
        FROM quiz_attempts qa 
        JOIN quizzes q ON qa.quiz_id = q.id 
        WHERE qa.user_id = ? 
        ORDER BY qa.completed_at DESC
    ''', (session['user_id'],)).fetchall()
    
    conn.close()
    
    # Convert attempts to a list of dictionaries and format dates
    attempts_list = []
    for attempt in attempts:
        attempt_dict = dict(attempt)
        # Format the completed_at date
        if attempt_dict['completed_at']:
            try:
                # Parse the date string and reformat it
                dt = datetime.strptime(attempt_dict['completed_at'], '%Y-%m-%d %H:%M:%S')
                attempt_dict['formatted_date'] = dt.strftime('%b %d, %Y at %I:%M %p')
            except (ValueError, TypeError):
                attempt_dict['formatted_date'] = "Unknown date"
        else:
            attempt_dict['formatted_date'] = "Unknown date"
        
        attempts_list.append(attempt_dict)
    
    return render_template("history.html", attempts=attempts_list)

# ---------------- VIEW ATTEMPT DETAILS ----------------
@routes.route("/attempt/<int:attempt_id>")
@login_required
def view_attempt(attempt_id):
    conn = get_db_connection()
    
    # Get the quiz attempt details
    attempt = conn.execute('''
        SELECT qa.*, q.title, q.context 
        FROM quiz_attempts qa 
        JOIN quizzes q ON qa.quiz_id = q.id 
        WHERE qa.id = ? AND qa.user_id = ?
    ''', (attempt_id, session['user_id'])).fetchone()
    
    if not attempt:
        flash("Quiz attempt not found", "error")
        return redirect(url_for("routes.history"))
    
    # Parse the answers
    try:
        user_answers = json.loads(attempt['answers'])
    except:
        user_answers = []
    
    # Get the questions for this quiz
    questions = conn.execute(
        'SELECT * FROM questions WHERE quiz_id = ?', (attempt['quiz_id'],)
    ).fetchall()
    
    conn.close()
    
    # Prepare data for display
    questions_data = []
    for i, question in enumerate(questions):
        question_dict = dict(question)
        question_dict['user_answer'] = user_answers[i] if i < len(user_answers) else "Not answered"
        question_dict['is_correct'] = question_dict['user_answer'] == question_dict['answer']
        
        # Parse options
        try:
            question_dict['options'] = json.loads(question_dict['options'])
        except:
            question_dict['options'] = []
        
        questions_data.append(question_dict)
    
    # Calculate score percentage
    score_percentage = (attempt['score'] / attempt['total_questions']) * 100 if attempt['total_questions'] > 0 else 0
    
    return render_template(
        "attempt_details.html", 
        attempt=attempt, 
        questions=questions_data,
        score_percentage=score_percentage
    )

# ---------------- PROGRESS PAGE ----------------
@routes.route("/progress")
@login_required
def progress():
    conn = get_db_connection()
    
    # Get user stats
    user_id = session['user_id']
    
    # Quizzes completed
    quizzes_completed = conn.execute(
        'SELECT COUNT(*) FROM quiz_attempts WHERE user_id = ?', (user_id,)
    ).fetchone()[0]
    
    # Average score
    avg_score_result = conn.execute(
        'SELECT AVG(score * 100.0 / total_questions) FROM quiz_attempts WHERE user_id = ?', (user_id,)
    ).fetchone()[0]
    average_score = round(avg_score_result, 1) if avg_score_result else 0
    
    # Quizzes created
    quizzes_created = conn.execute(
        'SELECT COUNT(*) FROM quizzes WHERE user_id = ?', (user_id,)
    ).fetchone()[0]
    
    # Recent quizzes (last 5 attempts)
    recent_attempts = conn.execute('''
        SELECT q.title, qa.score, qa.total_questions, qa.completed_at 
        FROM quiz_attempts qa 
        JOIN quizzes q ON qa.quiz_id = q.id 
        WHERE qa.user_id = ? 
        ORDER BY qa.completed_at DESC 
        LIMIT 5
    ''', (user_id,)).fetchall()
    
    # Format recent quizzes for display
    recent_quizzes = []
    for attempt in recent_attempts:
        recent_quizzes.append({
            'title': attempt['title'],
            'score': round((attempt['score'] / attempt['total_questions']) * 100),
            'date': attempt['completed_at'][:10]  # Just the date part
        })
    
    # Calculate progress percentage (quizzes completed vs goal of 20 quizzes)
    quizzes_goal = 20
    progress_percentage = min(100, round((quizzes_completed / quizzes_goal) * 100))
    
    # Calculate days streak (placeholder implementation)
    # In a real app, you would calculate this based on consecutive days with quiz activity
    days_streak = random.randint(1, 10)  # Placeholder
    
    # Check achievements
    achievements = {
        'first_quiz': quizzes_completed >= 1,
        'quiz_master': quizzes_completed >= 10,
        'perfect_score': any(attempt['score'] == attempt['total_questions'] for attempt in recent_attempts),
        'streak_7': days_streak >= 7,
        'streak_30': days_streak >= 30,
        'quick_learner': average_score >= 80,
        'earned': 0,
        'total': 6
    }
    
    # Count earned achievements
    achievements['earned'] = sum(1 for k, v in achievements.items() if v and k not in ['earned', 'total'])
    
    conn.close()
    
    # Prepare stats for template
    stats = {
        'quizzes_completed': quizzes_completed,
        'average_score': average_score,
        'quizzes_created': quizzes_created,
        'progress_percentage': progress_percentage,
        'quizzes_goal': quizzes_goal,
        'days_streak': days_streak,
        'achievements': achievements
    }
    
    return render_template("progress.html", user=session['user'], stats=stats, recent_quizzes=recent_quizzes)

# ---------------- SETTINGS PAGE ----------------
@routes.route("/settings")
@login_required
def settings():
    conn = get_db_connection()
    user = conn.execute(
        'SELECT * FROM users WHERE id = ?', (session['user_id'],)
    ).fetchone()
    conn.close()
    
    # Parse user preferences
    preferences = json.loads(user['preferences'])
    current_theme = preferences.get('theme', 'light')
    
    return render_template("settings.html", user=user, current_theme=current_theme)

# ---------------- UPDATE PROFILE ----------------
@routes.route("/update_profile", methods=["POST"])
@login_required
def update_profile():
    new_username = request.form.get("username")
    
    if not new_username or not is_valid_username(new_username):
        flash("Invalid username", "error")
        return redirect(url_for("routes.settings"))
    
    conn = get_db_connection()
    
    # Check if username is already taken by another user
    existing_user = conn.execute(
        'SELECT * FROM users WHERE username = ? AND id != ?', 
        (new_username, session['user_id'])
    ).fetchone()
    
    if existing_user:
        flash("Username already taken", "error")
        conn.close()
        return redirect(url_for("routes.settings"))
    
    # Update username
    conn.execute(
        'UPDATE users SET username = ? WHERE id = ?',
        (new_username, session['user_id'])
    )
    conn.commit()
    conn.close()
    
    session['user'] = new_username
    flash("Profile updated successfully", "success")
    return redirect(url_for("routes.settings"))

# ---------------- CHANGE PASSWORD ----------------
@routes.route("/change_password", methods=["POST"])
@login_required
def change_password():
    current_password = request.form.get("current_password")
    new_password = request.form.get("new_password")
    confirm_password = request.form.get("confirm_password")
    
    if not current_password or not new_password or not confirm_password:
        flash("All fields are required", "error")
        return redirect(url_for("routes.settings"))
    
    conn = get_db_connection()
    user = conn.execute(
        'SELECT * FROM users WHERE id = ?', (session['user_id'],)
    ).fetchone()
    
    if not check_password_hash(user['password_hash'], current_password):
        flash("Current password is incorrect", "error")
        conn.close()
        return redirect(url_for("routes.settings"))
    
    if new_password != confirm_password:
        flash("New passwords do not match", "error")
        conn.close()
        return redirect(url_for("routes.settings"))
    
    if not is_valid_password(new_password):
        flash("Password must be 8+ chars with uppercase, lowercase, digit, special char", "error")
        conn.close()
        return redirect(url_for("routes.settings"))
    
    # Update password
    new_password_hash = generate_password_hash(new_password)
    conn.execute(
        'UPDATE users SET password_hash = ? WHERE id = ?',
        (new_password_hash, session['user_id'])
    )
    conn.commit()
    conn.close()
    
    flash("Password updated successfully", "success")
    return redirect(url_for("routes.settings"))

# ---------------- UPDATE PREFERENCES ----------------
@routes.route("/update_preferences", methods=["POST"])
@login_required
def update_preferences():
    theme = request.form.get("theme", "light")
    
    if theme not in ["light", "dark", "purple"]:
        theme = "light"
    
    conn = get_db_connection()
    
    # Get current preferences
    user = conn.execute(
        'SELECT * FROM users WHERE id = ?', (session['user_id'],)
    ).fetchone()
    
    # Parse current preferences
    preferences = json.loads(user['preferences'])
    preferences['theme'] = theme
    
    # Update preferences
    conn.execute(
        'UPDATE users SET preferences = ? WHERE id = ?',
        (json.dumps(preferences), session['user_id'])
    )
    conn.commit()
    conn.close()
    
    flash("Preferences updated successfully", "success")
    return redirect(url_for("routes.settings"))

# ---------------- DELETE ACCOUNT ----------------
@routes.route("/delete_account", methods=["POST"])
@login_required
def delete_account():
    conn = get_db_connection()
    
    # Delete user and all their data
    conn.execute('DELETE FROM quiz_attempts WHERE user_id = ?', (session['user_id'],))
    conn.execute('DELETE FROM questions WHERE quiz_id IN (SELECT id FROM quizzes WHERE user_id = ?)', (session['user_id'],))
    conn.execute('DELETE FROM quizzes WHERE user_id = ?', (session['user_id'],))
    conn.execute('DELETE FROM users WHERE id = ?', (session['user_id'],))
    
    conn.commit()
    conn.close()
    
    # Clear session
    session.clear()
    
    flash("Your account has been deleted", "info")
    return redirect(url_for("routes.index"))

# ---------------- UPLOAD PAGE ----------------
@routes.route("/upload", methods=["GET"])
@login_required
def upload_page():
    return render_template("upload.html")

# ---------------- FILE UPLOAD ----------------
@routes.route("/upload", methods=["POST"])
@login_required
def upload_file():
    try:
        if "file" not in request.files:
            flash("No file provided", "error")
            return jsonify({"error": "No file provided"}), 400

        file = request.files["file"]
        filename = secure_filename(file.filename)
        context = ""

        if filename.endswith(".pdf"):
            pdf = PyPDF2.PdfReader(file)
            context = "".join(page.extract_text() or "" for page in pdf.pages)
        elif filename.endswith(".txt"):
            context = file.read().decode("utf-8")
        elif filename.endswith(".pptx"):
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        context += shape.text + "\n"
        else:
            flash("Unsupported file format", "error")
            return jsonify({"error": "Unsupported file format"}), 400

        flash("File uploaded successfully", "success")
        return jsonify({"context": context[:10000]})

    except Exception as e:
        flash(f"Error processing file: {str(e)}", "error")
        return jsonify({"error": f"Error processing file: {str(e)}"}), 500

# ---------------- QUIZ GENERATION ----------------
try:
    qg_pipeline = pipeline(
        "text2text-generation",
        model="valhalla/t5-base-e2e-qg",
        tokenizer=AutoTokenizer.from_pretrained("valhalla/t5-base-e2e-qg", use_fast=False)
    )
    qa_pipeline = pipeline(
        "question-answering",
        model="deepset/roberta-base-squad2",
        tokenizer=AutoTokenizer.from_pretrained("deepset/roberta-base-squad2", use_fast=False)
    )
except Exception as e:
    qg_pipeline = pipeline(
        "text2text-generation",
        model="valhalla/t5-small-e2e-qg"
    )
    qa_pipeline = pipeline(
        "question-answering",
        model="distilbert-base-cased-distilled-squad"
    )

def clean_question(question_text):
    """Clean up generated questions by removing tags and ensuring proper formatting"""
    if not question_text:
        return ""
    
    # Remove <sep> tags and any other HTML-like tags
    question_text = re.sub(r'<[^>]+>', '', question_text).strip()
    
    # Remove repetitive patterns (common issue with T5 model)
    question_text = re.sub(r'(\b\w+\b)(?:\s+\1)+', r'\1', question_text)
    
    # Extract only the first question if multiple are generated
    questions = re.split(r'\?+\s+', question_text)
    if questions:
        question_text = questions[0].strip()
        if not question_text.endswith('?'):
            question_text += '?'
    
    # Remove trailing special characters
    question_text = re.sub(r'[.,;:]+$', '', question_text)
    
    # Ensure it ends with a question mark
    if not question_text.endswith('?'):
        question_text = question_text.rstrip('.') + '?'
    
    # Capitalize first letter
    if question_text and question_text[0].islower():
        question_text = question_text[0].upper() + question_text[1:]
    
    # Limit length
    if len(question_text.split()) > 15:
        question_text = ' '.join(question_text.split()[:12]) + '...?'
    
    return question_text

def generate_distractors(answer, num_choices=3, context=None):
    answer_text = str(answer).strip()
    if not answer_text:
        return []

    
    # Remove trailing punctuation from answer for better matching
    answer_clean = re.sub(r'[.,;:]+$', '', answer_text)
    answer_lower = answer_clean.lower()
    
    distractors = set()
    
    # Use context for distractor generation
    if context:
        try:
            # Extract unique words from context (nouns and adjectives)
            words = re.findall(r'\b[A-Za-z]+\b', context)
            word_freq = Counter(words)
            
            # Get most frequent words that aren't the answer
            for word, count in word_freq.most_common(50):
                word_lower = word.lower()
                if (word_lower != answer_lower and 
                    len(word) > 3 and
                    word_lower not in {'the', 'and', 'for', 'with', 'this', 'that', 'which'}):
                    distractors.add(word)
                    if len(distractors) >= num_choices * 3:
                        break
        except Exception as e:
            print(f"⚠️ Context processing error: {e}")

    # Add topic-specific distractors
    topic_distractors = {
        "testing": ["development", "debugging", "quality assurance", "verification", "validation"],
        "date": ["yesterday", "tomorrow", "next week", "last month", "next meeting"],
        "developer": ["tester", "designer", "manager", "analyst", "architect"],
        "white box": ["black box", "gray box", "unit testing", "integration testing", "system testing"]
    }
    
    # Add relevant topic distractors
    for topic, terms in topic_distractors.items():
        if topic in answer_lower or any(topic in word for word in distractors):
            for term in terms:
                if term not in distractors:
                    distractors.add(term)
                if len(distractors) >= num_choices * 4:
                    break

    # Format and return
    final_distractors = []
    for distractor in list(distractors)[:num_choices * 2]:  # Get more options
        # Clean up the distractor
        distractor_clean = re.sub(r'[.,;:]+$', '', str(distractor)).strip()
        
        if distractor_clean and distractor_clean.lower() != answer_lower:
            if answer_text and answer_text[0].isupper():
                final_distractors.append(distractor_clean.capitalize())
            else:
                final_distractors.append(distractor_clean.lower())
    
    # Remove duplicates while preserving order
    seen = set()
    unique_distractors = []
    for distractor in final_distractors:
        if distractor.lower() not in seen:
            seen.add(distractor.lower())
            unique_distractors.append(distractor)
    
    return unique_distractors[:num_choices]

def generate_content_based_fallback(context, num_questions):
    try:
        words = re.findall(r'\b[a-zA-Z]+\b', context.lower())
        stop_words = {'the', 'and', 'for', 'with', 'this', 'that', 'which', 'what'}
        key_terms = [word for word in words if word not in stop_words and len(word) > 3]
        
        if not key_terms:
            key_terms = ["content", "topic", "subject", "information"]
        else:
            term_counts = Counter(key_terms)
            key_terms = [term for term, count in term_counts.most_common(8)]
        
        questions = []
        question_templates = [
            "What is the significance of {term}?",
            "How does {term} contribute?",
            "What role does {term} play?",
            "Why is {term} important?"
        ]
        
        for i in range(num_questions):
            term = key_terms[i % len(key_terms)]
            template = random.choice(question_templates)
            
            questions.append({
                "type": "MCQ",
                "question": template.format(term=term.capitalize()),
                "options": [
                    "Key element discussed",
                    "Important role", 
                    "Central to topic",
                    "Provides context"
                ],
                "answer": "Key element discussed"
            })
        
        return jsonify({
            "source": "CONTENT_BASED",
            "questions": questions
        })
        
    except Exception as e:
        return jsonify({
            "source": "BASIC",
            "questions": [
                {
                    "type": "MCQ",
                    "question": "What is the main topic discussed?",
                    "options": [
                        "The subject matter",
                        "Unrelated concepts", 
                        "General information",
                        "Various topics"
                    ],
                    "answer": "The subject matter"
                }
                for _ in range(num_questions)
            ]
        })

@routes.route("/generate", methods=["POST"])
@login_required
def generate_questions():
    try:
        data = request.get_json()
        num_questions = int(data.get("num_questions", 5))
        context = data.get("context", "").strip()

        print(f"📝 Context received: {len(context)} characters")
        if context:
            print(f"📝 First 200 chars: {context[:200]}...")

        if not context:
            return jsonify({"error": "Empty context"}), 400

        # Better chunking - use paragraphs instead of sentences
        chunks = re.split(r'\n\s*\n', context)  # Split by paragraphs
        chunks = [c.strip() for c in chunks if len(c.split()) > 10]  # Minimum 10 words
        
        # Fallback to sentences if no good paragraphs
        if not chunks:
            chunks = re.split(r'(?<=[.!?])\s+', context)
            chunks = [c.strip() for c in chunks if len(c.split()) > 8]

        print(f"📊 Number of chunks: {len(chunks)}")
        if chunks:
            print(f"📊 Sample chunk: {chunks[0][:100]}...")

        if not chunks:
            return generate_content_based_fallback(context, num_questions)

        response_questions = []
        used_questions = set()

        for attempt in range(num_questions * 5):  # More attempts
            if len(response_questions) >= num_questions:
                break

            chunk = random.choice(chunks)

            try:
                # Use a better prompt for question generation
                q_out = qg_pipeline(
                    f"Generate a question about: {chunk}",
                    max_new_tokens=50,  # Reduced from 64
                    num_return_sequences=1,
                    do_sample=True,
                    temperature=0.8,  # Slightly lower temperature
                    top_p=0.9,
                    repetition_penalty=1.2  # Reduce repetition
                )
                question_text = q_out[0]["generated_text"].strip()
                
                # Clean the question text
                question_text = clean_question(question_text)
                print(f"❓ Generated question: {question_text}")

                # Enhanced question validation
                if (not question_text.endswith("?") or
                    len(question_text.split()) < 4 or
                    len(question_text.split()) > 20 or
                    question_text.count('?') > 2 or  # Avoid multiple questions
                    any(word in question_text.lower() for word in ['example:', 'labels:', 'step:']) or
                    question_text in used_questions):
                    print("⚠️  Skipping: Invalid question format or duplicate")
                    continue

                used_questions.add(question_text)

                # Use the entire context for answer extraction for better results
                answer_result = qa_pipeline({
                    "question": question_text,
                    "context": context  # Use full context instead of just chunk
                })
                answer = answer_result.get("answer", "").strip()
                confidence = answer_result.get("score", 0.0)
                
                # Clean the answer
                answer = re.sub(r'^[^a-zA-Z0-9]*|[^a-zA-Z0-9]*$', '', answer)
                
                print(f"✅ Answer found: {answer} (confidence: {confidence:.2f})")

                # Better answer validation
                if (not answer or 
                    confidence < 0.1 or 
                    len(answer.split()) > 8 or  # Avoid very long answers
                    answer.lower() in {'yes', 'no', 'true', 'false'}):  # Avoid simple answers
                    print("⚠️  Skipping: No answer or low confidence")
                    continue

                # Generate distractors
                distractors = generate_distractors(answer, 3, context=context)
                print(f"🎯 Distractors: {distractors}")

                options = [answer] + [d for d in distractors if d and d.lower() != answer.lower()]
                options = list(dict.fromkeys(options))

                # Better fallback options
                while len(options) < 4:
                    # Use context words for fallback options
                    context_words = re.findall(r'\b[A-Z][a-z]+\b', context)
                    if context_words:
                        new_opt = random.choice(context_words)
                        if new_opt not in options and new_opt != answer:
                            options.append(new_opt)
                    else:
                        options.append(f"Option {chr(65 + len(options))}")

                random.shuffle(options)
                print(f"✅ Final options: {options}")

                response_questions.append({
                    "type": "MCQ",
                    "question": question_text,
                    "options": options[:4],
                    "answer": answer
                })

            except Exception as e:
                print(f"⚠️ Generation error: {e}")
                continue

        if not response_questions:
            print("⚠️  No questions generated, using fallback")
            return generate_content_based_fallback(context, num_questions)

        # Store the quiz in database - FIXED: Use cursor for lastrowid
        conn = get_db_connection()
        cursor = conn.cursor()
        
        # Save quiz
        cursor.execute(
            'INSERT INTO quizzes (user_id, title, context) VALUES (?, ?, ?)',
            (session['user_id'], 'Generated Quiz', context[:500])
        )
        quiz_id = cursor.lastrowid
        
        # Save questions
        for q in response_questions:
            cursor.execute(
                'INSERT INTO questions (quiz_id, question_text, options, answer) VALUES (?, ?, ?, ?)',
                (quiz_id, q['question'], json.dumps(q['options']), q['answer'])
            )
        
        conn.commit()
        conn.close()
        
        # Store the quiz in session for the interactive quiz
        session['current_quiz'] = response_questions
        session['current_quiz_id'] = quiz_id
        session.modified = True  # Ensure session changes are saved
        
        return jsonify({
            "source": "AI",
            "questions": response_questions,
            "message": f"Generated {len(response_questions)} questions",
            "redirect": url_for("routes.take_quiz")
        })

    except Exception as e:
        print(f"❌ Error in generate_questions: {e}")
        return jsonify({"error": str(e)}), 500

# ---------------- INTERACTIVE QUIZ PAGE ----------------
@routes.route("/take_quiz")
@login_required
def take_quiz():
    # Get the quiz data from session
    quiz_data = session.get('current_quiz', [])
    
    if not quiz_data:
        # If no quiz data, redirect to dashboard with a message
        flash("No quiz available. Please generate a quiz first.", "warning")
        return redirect(url_for("routes.dashboard"))
    
    # Ensure the quiz data is properly formatted
    for question in quiz_data:
        if 'options' not in question:
            question['options'] = []
        if 'answer' not in question:
            question['answer'] = ""
    
    return render_template("quiz.html", questions=quiz_data)

# ---------------- DOWNLOAD PDF ----------------
@routes.route("/download_pdf", methods=["POST"])
@login_required
def download_pdf():
    try:
        from fpdf import FPDF
        
        topic = request.form.get("topic", "Generated Quiz")
        questions_data = request.form.get("questions", "[]")
        
        try:
            questions = json.loads(questions_data)
        except Exception as e:
            flash("Invalid questions data", "error")
            return redirect(url_for("routes.dashboard"))

        # Create PDF
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        
        # Add title
        pdf.set_font("Arial", 'B', 16)
        pdf.cell(200, 10, txt=f"Quiz: {topic}", ln=True, align='C')
        pdf.ln(10)
        
        # Reset to regular font
        pdf.set_font("Arial", size=12)
        
        # Add questions
        for i, q in enumerate(questions, 1):
            # Question text
            pdf.set_font("Arial", 'B', 12)
            pdf.multi_cell(0, 10, f"{i}. {q.get('question', 'Question')}")
            pdf.set_font("Arial", size=12)
            
            # Options
            if isinstance(q.get('options'), list):
                for j, opt in enumerate(q['options']):
                    pdf.cell(10)  # Indent
                    pdf.multi_cell(0, 8, f"{chr(65 + j)}. {opt}")
            
            # Answer
            pdf.set_font("Arial", 'I', 12)
            pdf.multi_cell(0, 10, f"Answer: {q.get('answer', 'Not provided')}")
            pdf.set_font("Arial", size=12)
            
            pdf.ln(5)  # Space between questions

        # Create a temporary file
        import tempfile
        import os
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
            pdf.output(tmp.name)
            tmp_path = tmp.name

        # Send the file
        response = send_file(
            tmp_path,
            as_attachment=True,
            download_name=f"{topic.replace(' ', '_')}_Quiz.pdf",
            mimetype='application/pdf'
        )
        
        # Clean up the temporary file after sending
        @response.call_on_close
        def remove_temp_file():
            try:
                os.unlink(tmp_path)
            except:
                pass
                
        return response

    except Exception as e:
        print(f"PDF Generation Error: {str(e)}")
        flash("Failed to generate PDF", "error")
        return redirect(url_for("routes.dashboard"))

# ---------------- SAVE QUIZ ATTEMPT ----------------
@routes.route("/save_attempt", methods=["POST"])
@login_required
def save_attempt():
    try:
        data = request.get_json()
        score = data.get('score')
        total = data.get('total')
        answers = data.get('answers')
        quiz_id = session.get('current_quiz_id')
        
        print(f"Save attempt - User ID: {session['user_id']}, Quiz ID: {quiz_id}, Score: {score}/{total}")
        
        if not quiz_id:
            return jsonify({"error": "No quiz ID found"}), 400
        
        conn = get_db_connection()
        conn.execute(
            'INSERT INTO quiz_attempts (user_id, quiz_id, score, total_questions, answers) VALUES (?, ?, ?, ?, ?)',
            (session['user_id'], quiz_id, score, total, json.dumps(answers))
        )
        conn.commit()
        conn.close()
        
        print("Quiz attempt saved successfully")
        return jsonify({"success": True})
    
    except Exception as e:
        print(f"Error saving quiz attempt: {e}")
        return jsonify({"error": str(e)}), 500

# ---------------- LOGOUT ----------------
@routes.route("/logout")
def logout():
    session.clear()
    return redirect(url_for("routes.login"))