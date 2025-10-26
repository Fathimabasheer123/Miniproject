from flask import Blueprint, render_template, request, jsonify, session, flash, send_file, url_for, redirect
import PyPDF2
import json
import os
import re
import tempfile
import io
from pptx import Presentation
from docx import Document
from fpdf import FPDF
from datetime import datetime
from werkzeug.utils import secure_filename
from reportlab.lib.pagesizes import letter, A4
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.units import inch
from reportlab.lib import colors

# OCR imports with path configuration
try:
    import pytesseract
    from PIL import Image
    
    # Set Tesseract path explicitly
    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
    
    # Test if Tesseract is accessible
    try:
        version = pytesseract.get_tesseract_version()
        OCR_AVAILABLE = True
        print(f"✅ OCR configured successfully! Tesseract version: {version}")
    except Exception as e:
        print(f"❌ Tesseract not accessible: {e}")
        OCR_AVAILABLE = False
        
except ImportError as e:
    OCR_AVAILABLE = False
    print(f"⚠️ OCR dependencies not installed: {e}")

# PDF processing imports
try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
    print("✅ pdfplumber available for enhanced PDF processing")
except ImportError:
    PDFPLUMBER_AVAILABLE = False
    print("⚠️ pdfplumber not available, using PyPDF2 only")

# pdf2image for scanned PDFs
try:
    from pdf2image import convert_from_bytes
    PDF2IMAGE_AVAILABLE = True
    print("✅ pdf2image available for scanned PDF processing")
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    print("⚠️ pdf2image not available")

from app.models.database import get_db_connection, login_required
from app.services.gemini_api import gemini_api

quiz_bp = Blueprint('quiz', __name__)

def extract_text_from_scanned_pdf(pdf_file):
    """Extract text from scanned PDF using OCR"""
    try:
        # Method 1: Use pdf2image + pytesseract
        if PDF2IMAGE_AVAILABLE:
            try:
                images = convert_from_bytes(pdf_file.read())
                text = ""
                for i, image in enumerate(images):
                    page_text = pytesseract.image_to_string(image)
                    text += f"Page {i+1}:\n{page_text}\n\n"
                print(f"✅ Extracted {len(text)} characters from scanned PDF via pdf2image OCR")
                return text
            except Exception as e:
                print(f"pdf2image OCR failed: {e}")
        
        # Method 2: Fallback - provide user instructions
        print("⚠️ Advanced scanned PDF processing not fully available")
        return ""
        
    except Exception as e:
        print(f"Scanned PDF extraction error: {e}")
        return ""

def enhanced_pdf_extraction(pdf_file):
    """Try multiple methods to extract text from PDF"""
    best_text = ""
    best_method = ""
    
    # Reset file pointer
    pdf_file.seek(0)
    
    # Method 1: Try pdfplumber (best for text PDFs)
    if PDFPLUMBER_AVAILABLE:
        try:
            with pdfplumber.open(pdf_file) as pdf:
                text = "".join(page.extract_text() or "" for page in pdf.pages)
                if len(text.strip()) > len(best_text.strip()):
                    best_text = text
                    best_method = "pdfplumber"
                    print(f"📄 pdfplumber extracted {len(text)} characters")
        except Exception as e:
            print(f"pdfplumber failed: {e}")
    
    # Method 2: Try PyPDF2 (fallback)
    pdf_file.seek(0)
    try:
        reader = PyPDF2.PdfReader(pdf_file)
        text = "".join(page.extract_text() or "" for page in reader.pages)
        if len(text.strip()) > len(best_text.strip()):
            best_text = text
            best_method = "PyPDF2"
            print(f"📄 PyPDF2 extracted {len(text)} characters")
    except Exception as e:
        print(f"PyPDF2 failed: {e}")
    
    if best_text:
        print(f"✅ Best PDF extraction: {best_method} with {len(best_text)} characters")
    
    return best_text, best_method

def extract_quiz_title(context, custom_title=None):
    """Extract a meaningful quiz title from the context or use custom title"""
    # PRIORITY 1: Use custom title if provided
    if custom_title and custom_title.strip():
        title = custom_title.strip()
        if len(title) > 60:
            title = title[:57] + "..."
        return title
    
    # PRIORITY 2: Try to extract from context
    if not context:
        return "Generated Quiz"
    
    # Try to extract first meaningful sentence
    sentences = re.split(r'[.!?]', context.strip())
    for sentence in sentences:
        sentence = sentence.strip()
        if len(sentence) > 20 and len(sentence.split()) >= 3:
            # Clean up the sentence for title
            title = sentence[:60].strip()  # Limit length
            if not title.endswith(('.', '!', '?')):
                title += '...'
            return title
    
    # Fallback: use first 4 words
    words = context.strip().split()[:4]
    if words:
        return ' '.join(words) + '...'
    
    return "Study Material Quiz"

@quiz_bp.route("/dashboard")
@login_required
def dashboard():
    """User dashboard"""
    conn = get_db_connection()
    uid = session['user_id']
    
    # Get basic stats - USE SAME LOGIC AS HISTORY PAGE
    quizzes_completed = conn.execute(
        'SELECT COUNT(DISTINCT quiz_id) FROM quiz_attempts WHERE user_id=?', (uid,)
    ).fetchone()[0]

# FIX: Use only latest attempts per quiz (same as history page)
    avg_score_result = conn.execute('''
        SELECT AVG(score * 100.0 / total_questions) 
        FROM quiz_attempts 
        WHERE id IN (
            SELECT MAX(id) 
            FROM quiz_attempts 
            WHERE user_id = ? 
            GROUP BY quiz_id
        )
    ''', (uid,)).fetchone()[0]
    average_score = round(avg_score_result, 1) if avg_score_result else 0

    quizzes_created = conn.execute(
        'SELECT COUNT(*) FROM quizzes WHERE user_id=?', (uid,)
    ).fetchone()[0]

# FIX: Also update total_time to use same logic
    total_time = conn.execute('''
        SELECT COALESCE(SUM(time_spent), 0) 
        FROM quiz_attempts 
        WHERE id IN (
            SELECT MAX(id) 
            FROM quiz_attempts 
            WHERE user_id = ? 
            GROUP BY quiz_id
        )
    ''', (uid,)).fetchone()[0]
    
    streak_result = conn.execute('''
        SELECT COUNT(DISTINCT DATE(completed_at)) as streak
        FROM quiz_attempts 
        WHERE user_id = ? AND completed_at >= date('now', '-30 days')
        ORDER BY completed_at DESC
    ''', (uid,)).fetchone()
    days_streak = streak_result['streak'] if streak_result else 0
    
    # Get recent quizzes with proper formatting - Only show latest attempt per quiz
    recent = conn.execute('''
        SELECT q.title, q.context, qa.score, qa.total_questions, qa.completed_at, q.difficulty
        FROM quiz_attempts qa
        JOIN quizzes q ON qa.quiz_id = q.id
        WHERE qa.user_id = ?
        AND qa.id IN (
            SELECT MAX(id) FROM quiz_attempts 
            WHERE user_id = ? 
            GROUP BY quiz_id
        )
        ORDER BY qa.completed_at DESC
        LIMIT 5
    ''', (uid, uid)).fetchall()
    conn.close()

    # Format recent quizzes with proper titles
    recent_quizzes = []
    for r in recent:
        # Use the actual title from database (now includes custom titles)
        title = r['title'] if r['title'] else 'Study Quiz'
        
        # Format date properly
        completed_date = r['completed_at']
        if completed_date:
            if 'T' in completed_date:
                dt = datetime.fromisoformat(completed_date.replace('Z', '+00:00'))
            else:
                dt = datetime.strptime(completed_date, '%Y-%m-%d %H:%M:%S')
            formatted_date = dt.strftime('%b %d, %Y')
        else:
            formatted_date = 'Unknown date'
        
        recent_quizzes.append({
            'title': title,
            'score': round((r['score'] / r['total_questions']) * 100) if r['total_questions'] > 0 else 0,
            'date': formatted_date,
            'difficulty': r['difficulty'] if r['difficulty'] else 'medium'
        })
           
   # Calculate progress based on actual quizzes created
    quizzes_goal = quizzes_created if quizzes_created > 0 else 1
    progress_percentage = min(100, round((quizzes_completed / quizzes_goal) * 100))

    stats = {
    'quizzes_completed': quizzes_completed,
    'average_score': average_score,
    'quizzes_created': quizzes_created,
    'total_time_spent': total_time // 60,
    'progress_percentage': progress_percentage,
    'quizzes_goal': quizzes_goal,  # Now shows actual available quizzes
    'days_streak': days_streak
}

    achievements = {
        'first_quiz': quizzes_completed >= 1,
        'quiz_master': quizzes_completed >= 10,
        'perfect_score': any(r['score'] == r['total_questions'] for r in recent),
        'streak_7': days_streak >= 7,
        'gemini_expert': quizzes_created >= 5,
        'difficulty_master': any(r['difficulty'] == 'hard' for r in recent if r['difficulty']),
        'earned': 0,
        'total': 6
    }
    achievements['earned'] = sum(
        1 for k, v in achievements.items()
        if v and k not in ['earned', 'total']
    )
    stats['achievements'] = achievements

    return render_template(
        "dashboard.html",
        user=session['user'],
        stats=stats,
        recent_quizzes=recent_quizzes
    )

@quiz_bp.route("/upload", methods=["GET", "POST"])
@login_required
def upload_file():
    """Handle file upload and text extraction"""
    if request.method == "GET":
        return render_template("upload.html")
    
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400
    
    f = request.files["file"]
    if f.filename == '':
        return jsonify({"error": "No file selected"}), 400
    
    filename = secure_filename(f.filename)
    context = ""
    
    try:
        # Store original filename for potential title use
        original_filename = os.path.splitext(filename)[0]  # Remove extension
        
        if filename.endswith(".pdf"):
            # Use enhanced PDF extraction
            context, method = enhanced_pdf_extraction(f)
            
            # If little/no text found, it's probably a scanned PDF
            if len(context.strip()) < 100:
                print("🔍 PDF appears to be scanned/handwritten - attempting OCR...")
                
                if not OCR_AVAILABLE:
                    return jsonify({
                        "error": "Scanned PDF detected but OCR not available. Please upload individual image files instead."
                    }), 400
                
                if not PDF2IMAGE_AVAILABLE:
                    return jsonify({
                        "error": "Scanned PDF detected. For best results, please save each page as an image (PNG/JPG) and upload separately.",
                        "suggestion": "Use 'Save as Image' in your PDF viewer or take screenshots of each page"
                    }), 400
                
                # Reset file pointer and try OCR
                f.seek(0)
                scanned_text = extract_text_from_scanned_pdf(f)
                
                if scanned_text and len(scanned_text.strip()) > 50:
                    context = scanned_text
                    print(f"✅ Successfully extracted {len(context)} characters from scanned PDF")
                else:
                    return jsonify({
                        "error": "Could not extract sufficient text from scanned PDF.",
                        "suggestions": [
                            "Ensure the PDF contains clear, readable text",
                            "Try uploading individual page images instead",
                            "Use a higher quality scan with better contrast"
                        ]
                    }), 400
            
        elif filename.endswith(".txt"):
            context = f.read().decode("utf-8")
            print(f"📝 Text file extracted {len(context)} characters")
            
        elif filename.endswith(".pptx"):
            prs = Presentation(f)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        context += shape.text + "\n"
            print(f"📊 PowerPoint extracted {len(context)} characters")
        
        elif filename.endswith(".docx"):
            doc = Document(f)
            context = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            print(f"📄 Word document extracted {len(context)} characters")
            
        elif filename.lower().endswith(('.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif')):
            if not OCR_AVAILABLE:
                return jsonify({"error": "OCR not available. Please use text-based files (PDF, TXT, DOCX, PPTX)."}), 400
                
            try:
                f.seek(0)
                image_data = f.read()
                print(f"🖼️ Processing image: {filename}, Size: {len(image_data)} bytes")
                
                image = Image.open(io.BytesIO(image_data))
                context = pytesseract.image_to_string(image)
                
                print(f"📝 OCR extracted {len(context)} characters from image")
                
                if not context.strip():
                    return jsonify({"error": "No text detected in image. Please ensure the image is clear and contains readable text."}), 400
                    
            except Exception as ocr_error:
                print(f"❌ OCR Error: {str(ocr_error)}")
                return jsonify({"error": f"OCR processing failed: {str(ocr_error)}"}), 400
                
        else:
            supported_formats = "PDF, TXT, PPTX, DOCX"
            if OCR_AVAILABLE:
                supported_formats += ", PNG, JPG, JPEG"
            return jsonify({"error": f"Unsupported file format. Please upload {supported_formats}."}), 400
        
        # Clean and limit context
        original_length = len(context)
        context = re.sub(r'\s+', ' ', context).strip()
        context = context[:10000]
        
        print(f"✅ Final context: {len(context)} characters (from {original_length} original)")
        
        if len(context) < 100:
            return jsonify({"error": "Not enough text extracted from file. Please try a different file or ensure it contains sufficient readable content."}), 400
        
        # Return both context and original filename for title generation
        return jsonify({
            "context": context,
            "filename": original_filename
        })
        
    except Exception as e:
        print(f"❌ Upload Error: {e}")
        return jsonify({"error": f"Error processing file: {str(e)}"}), 500

@quiz_bp.route("/generate", methods=["POST"])
@login_required
def generate_questions():
    """Generate questions using Gemini API"""
    try:
        data = request.get_json()
        num_questions = int(data.get("num_questions", 10))
        difficulty = data.get("difficulty", "medium")
        context = data.get("context", "").strip()
        custom_title = data.get("title", "").strip()  # Get custom title from frontend
        
        # Reduce context size to prevent token limits
        context = context[:3000]
        
        print(f"🧠 Generating {num_questions} {difficulty} questions from {len(context)} chars of context")
        print(f"📝 Custom title: {custom_title}")
        
        if not context:
            return jsonify({"error": "Empty context"}), 400

        if not gemini_api.is_available():
            return jsonify({"error": "Gemini AI service is not available. Please check your API key configuration."}), 500

        gemini_questions = gemini_api.generate_questions(context, num_questions, difficulty)
        
        if not gemini_questions:
            return jsonify({"error": "Failed to generate questions with Gemini AI. Please try again with different content."}), 500

        question_types = list(set([q.get('type', 'MCQ') for q in gemini_questions]))
        print(f"✅ Generated {len(gemini_questions)} questions of types: {question_types}")

        # Extract meaningful title - NOW USING CUSTOM TITLE
        quiz_title = extract_quiz_title(context, custom_title)
        print(f"🎯 Final quiz title: {quiz_title}")
        
        conn = get_db_connection()
        cur = conn.cursor()
        cur.execute(
            'INSERT INTO quizzes (user_id, title, context, difficulty, question_types) VALUES (?, ?, ?, ?, ?)',
            (session['user_id'], quiz_title, context[:500], difficulty, json.dumps(question_types))
        )
        quiz_id = cur.lastrowid
        
        for q in gemini_questions:
            cur.execute(
                '''INSERT INTO questions (quiz_id, question_text, options, answer, question_type, explanation, difficulty) 
                   VALUES (?, ?, ?, ?, ?, ?, ?)''',
                (quiz_id, q['question'], json.dumps(q['options']), q['answer'], 
                 q.get('type', 'MCQ'), q.get('explanation', ''), q.get('difficulty', difficulty))
            )
        conn.commit()
        conn.close()

        session['current_quiz'] = gemini_questions
        session['current_quiz_id'] = quiz_id
        session['quiz_title'] = quiz_title  # Store title in session for display

        return jsonify({
            "source": "Gemini AI",
            "questions": gemini_questions,
            "title": quiz_title,  # Send title back to frontend
            "message": f"Successfully generated {len(gemini_questions)} questions with explanations",
            "redirect": url_for("quiz.take_quiz")
        })

    except Exception as e:
        print(f"❌ Question generation failed: {e}")
        return jsonify({"error": f"Question generation failed: {str(e)}"}), 500

@quiz_bp.route("/take_quiz")
@login_required
def take_quiz():
    """Display the generated quiz"""
    quiz = session.get('current_quiz', [])
    quiz_title = session.get('quiz_title', 'Generated Quiz')
    
    if not quiz:
        flash("No quiz available. Please generate a quiz first.", "warning")
        return redirect(url_for("quiz.dashboard"))
    
    return render_template("quiz.html", questions=quiz, quiz_title=quiz_title)

@quiz_bp.route("/save_attempt", methods=["POST"])
@login_required
def save_attempt():
    """Save quiz attempt results - Update existing attempt for retakes"""
    data = request.get_json()
    quiz_id = session.get('current_quiz_id')
    user_id = session['user_id']
    
    if not quiz_id:
        return jsonify({"error": "No quiz id found"}), 400
    
    # Get the original questions to understand question types
    conn = get_db_connection()
    questions = conn.execute(
        'SELECT question_type, answer FROM questions WHERE quiz_id = ?', (quiz_id,)
    ).fetchall()
    conn.close()
    
    # Normalize answers based on question type
    normalized_answers = []
    user_answers = data.get('answers', [])
    
    for i, answer in enumerate(user_answers):
        if i < len(questions):
            question_type = questions[i]['question_type']
            correct_answer = questions[i]['answer']
            
            if question_type == 'TrueFalse':
                # For TrueFalse questions, convert letter back to text
                if answer == 'A':
                    normalized_answers.append('True')
                elif answer == 'B':
                    normalized_answers.append('False')
                else:
                    normalized_answers.append(answer or '')
            else:
                # For other question types, use letter format
                normalized = answer.strip().upper()[0] if answer and answer.strip() else ''
                normalized_answers.append(normalized)
        else:
            normalized_answers.append(answer or '')
    
    conn = get_db_connection()
    
    # Check if this is a retake
    is_retake = session.get('is_retake', False)
    
    print(f"💾 Saving attempt - Quiz ID: {quiz_id}, User ID: {user_id}, Is Retake: {is_retake}")
    
    if is_retake:
        # For retakes, update the most recent attempt
        latest_attempt = conn.execute('''
            SELECT id FROM quiz_attempts 
            WHERE user_id = ? AND quiz_id = ? 
            ORDER BY completed_at DESC LIMIT 1
        ''', (user_id, quiz_id)).fetchone()
        
        if latest_attempt:
            conn.execute('''
                UPDATE quiz_attempts 
                SET score = ?, total_questions = ?, answers = ?, time_spent = ?, completed_at = CURRENT_TIMESTAMP
                WHERE id = ?
            ''', (data['score'], data['total'], json.dumps(normalized_answers), 
                  data.get('time_spent', 0), latest_attempt['id']))
            attempt_id = latest_attempt['id']
            print(f"✅ Updated existing attempt {attempt_id} for retake")
        else:
            attempt_id = create_new_attempt(conn, user_id, quiz_id, data, normalized_answers)
    else:
        attempt_id = create_new_attempt(conn, user_id, quiz_id, data, normalized_answers)
    
    conn.commit()
    conn.close()
    
    # Clear session data
    session.pop('current_quiz', None)
    session.pop('current_quiz_id', None)
    session.pop('is_retake', None)
    session.pop('quiz_title', None)
    
    return jsonify({"success": True, "attempt_id": attempt_id})

def create_new_attempt(conn, user_id, quiz_id, data, normalized_answers):
    """Helper function to create a new quiz attempt"""
    quiz = conn.execute(
        'SELECT difficulty, question_types FROM quizzes WHERE id = ?', (quiz_id,)
    ).fetchone()
    
    cursor = conn.execute(
        '''INSERT INTO quiz_attempts (user_id, quiz_id, score, total_questions, answers, time_spent, difficulty, question_types) 
        VALUES (?, ?, ?, ?, ?, ?, ?, ?)''',
        (user_id, quiz_id, data['score'], data['total'], 
         json.dumps(normalized_answers), data.get('time_spent', 0),
         quiz['difficulty'] if quiz else 'medium', 
         quiz['question_types'] if quiz else '[]')
    )
    attempt_id = cursor.lastrowid
    print(f"✅ Created new quiz attempt {attempt_id}")
    return attempt_id

@quiz_bp.route("/download_pdf", methods=["POST"])
@login_required
def download_pdf():
    """Download quiz as PDF using reportlab for better formatting"""
    try:
        topic = request.form.get("topic", "Generated Quiz")
        questions_data = request.form.get("questions", "[]")
        
        try:
            questions = json.loads(questions_data)
        except Exception as e:
            print(f"❌ PDF Error: Invalid questions data - {e}")
            flash("Invalid questions data for PDF generation", "error")
            return redirect(url_for("quiz.dashboard"))

        # Use reportlab for PDF generation
        buffer = generate_pdf_with_reportlab_simple(questions, topic)
        
        # Create filename
        filename = f"QuizGen_{topic.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.pdf"
        
        response = send_file(
            buffer,
            as_attachment=True,
            download_name=filename,
            mimetype='application/pdf',
           
        )
                
        print(f"✅ PDF generated successfully using reportlab: {topic}")
        return response

    except Exception as e:
        print(f"❌ PDF Generation Error: {str(e)}")
        flash("Failed to generate PDF. Please try again.", "error")
        return redirect(url_for("quiz.dashboard"))

def generate_pdf_with_reportlab_simple(questions, topic):
    """Simple PDF generation using reportlab canvas"""
    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter
    
    # Set starting position
    y_position = height - 72  # Start from top with 1-inch margin
    
    # Title
    c.setFont("Helvetica-Bold", 16)
    c.setFillColorRGB(0.29, 0.04, 0.46)  # Purple color
    title = f"Quiz: {topic}"
    c.drawString(72, y_position, title[:80] + "..." if len(title) > 80 else title)
    y_position -= 30
    
    # Quiz info
    c.setFont("Helvetica", 10)
    c.setFillColorRGB(0, 0, 0)  # Black color
    c.drawString(72, y_position, f"Total Questions: {len(questions)}")
    y_position -= 15
    c.drawString(72, y_position, f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    y_position -= 30
    
    # Process questions
    for i, q in enumerate(questions, 1):
        # Check if we need a new page
        if y_position < 100:
            c.showPage()
            y_position = height - 72
        
        # Question number and text
        c.setFont("Helvetica-Bold", 12)
        question_text = f"{i}. {q.get('question', 'Question')}"
        
        # Split long questions into multiple lines
        words = question_text.split()
        lines = []
        current_line = []
        
        for word in words:
            test_line = ' '.join(current_line + [word])
            if c.stringWidth(test_line, "Helvetica-Bold", 12) < (width - 144):  # Account for margins
                current_line.append(word)
            else:
                if current_line:
                    lines.append(' '.join(current_line))
                current_line = [word]
        if current_line:
            lines.append(' '.join(current_line))
        
        # Draw question lines
        for line in lines:
            if y_position < 100:
                c.showPage()
                y_position = height - 72
            c.drawString(72, y_position, line)
            y_position -= 15
        
        # Options
        c.setFont("Helvetica", 10)
        options = q.get('options', [])
        for j, opt in enumerate(options):
            if y_position < 80:
                c.showPage()
                y_position = height - 72
            option_text = f"{chr(65 + j)}. {opt}"
            # Truncate very long options
            if len(option_text) > 100:
                option_text = option_text[:97] + "..."
            c.drawString(85, y_position, option_text)
            y_position -= 12
        
        # Answer
        if y_position < 80:
            c.showPage()
            y_position = height - 72
        c.setFont("Helvetica-Oblique", 10)
        answer = q.get('answer', 'Not provided')
        c.drawString(72, y_position, f"Answer: {answer}")
        y_position -= 15
        
        # Explanation
        explanation = q.get('explanation', '')
        if explanation:
            if y_position < 80:
                c.showPage()
                y_position = height - 72
            c.setFont("Helvetica", 9)
            expl_text = f"Explanation: {explanation}"
            # Split long explanations
            expl_words = expl_text.split()
            expl_lines = []
            current_expl_line = []
            
            for word in expl_words:
                test_line = ' '.join(current_expl_line + [word])
                if c.stringWidth(test_line, "Helvetica", 9) < (width - 144):
                    current_expl_line.append(word)
                else:
                    if current_expl_line:
                        expl_lines.append(' '.join(current_expl_line))
                    current_expl_line = [word]
            if current_expl_line:
                expl_lines.append(' '.join(current_expl_line))
            
            for line in expl_lines:
                if y_position < 80:
                    c.showPage()
                    y_position = height - 72
                c.drawString(72, y_position, line)
                y_position -= 10
        
        y_position -= 10  # Space between questions
    
    c.save()
    buffer.seek(0)
    return buffer